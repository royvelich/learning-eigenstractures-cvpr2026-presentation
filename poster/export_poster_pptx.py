"""
Export the poster PDF to a single-slide PowerPoint file.

Strategy: render the PDF to a high-res PNG, then embed that PNG into a
PPTX slide whose dimensions match the physical poster (213 cm x 107 cm).

Usage:
    gen/.venv/Scripts/python.exe poster/export_poster_pptx.py
"""
from pathlib import Path
import subprocess
import shutil
import sys

from pptx import Presentation
from pptx.util import Cm, Emu

ROOT = Path(__file__).resolve().parent
PDF = ROOT / "poster.pdf"
PNG_PREFIX = ROOT / "poster_export"
PPTX = ROOT / "poster.pptx"

# Physical poster size = 213 cm x 107 cm, but PowerPoint caps slide
# dimensions at 56 inches (~142.2 cm). We scale down to 142 x 71.3 cm so
# the slide fits PowerPoint's limit while preserving the 213/107 aspect.
POSTER_ASPECT = 213 / 107
POSTER_W_CM = 142.0
POSTER_H_CM = POSTER_W_CM / POSTER_ASPECT  # ~71.3 cm

# Render the PDF to PNG via MiKTeX's pdftoppm. 150 DPI yields ~12.6 K x
# ~6.3 K px at this poster size — a good balance of clarity and file size.
DPI = 150

PDFTOPPM = shutil.which("pdftoppm") or (
    r"C:\Users\Roy\AppData\Local\Programs\MiKTeX\miktex\bin\x64\pdftoppm.exe"
)

print(f"[..] rendering {PDF.name} at {DPI} dpi via {Path(PDFTOPPM).name}")
# pdftoppm writes <prefix>-1.png for a single-page PDF
for stale in ROOT.glob("poster_export-*.png"):
    stale.unlink()
subprocess.run(
    [PDFTOPPM, "-r", str(DPI), "-png", str(PDF), str(PNG_PREFIX)],
    check=True,
)
png_path = next(ROOT.glob("poster_export-*.png"))
print(f"[ok] wrote {png_path.name} ({png_path.stat().st_size/1024/1024:.1f} MB)")

# Build a PPTX with a single slide sized to the physical poster.
prs = Presentation()
prs.slide_width = Cm(POSTER_W_CM)
prs.slide_height = Cm(POSTER_H_CM)

# Blank layout (layout index 6 is "Blank" in the default template).
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Insert the rendered poster at (0, 0) filling the full slide.
slide.shapes.add_picture(
    str(png_path),
    left=Emu(0),
    top=Emu(0),
    width=prs.slide_width,
    height=prs.slide_height,
)

prs.save(str(PPTX))
print(f"[ok] wrote {PPTX.name} ({PPTX.stat().st_size/1024/1024:.1f} MB)")
print(f"     slide size: {POSTER_W_CM} cm x {POSTER_H_CM} cm")
